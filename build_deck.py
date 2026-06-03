#!/usr/bin/env python3
"""
Build the Orient keynote pitch deck as a .pptx (opens natively in Apple Keynote).

18 slides, human-loop-first, styled to match the Orient site:
  - palette derived from the site's OKLCH tokens (orient.css)
  - Spectral (display serif) / IBM Plex Sans (UI) / IBM Plex Mono (labels)
  - native vector logomark + the real product screenshots

Run:  python3 build_deck.py
Out:  Orient-Keynote.pptx
"""

import math
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
SHOTS = os.path.join(HERE, "screenshots")

# ---------------------------------------------------------------------------
# Colour — convert the site's OKLCH tokens to sRGB so the deck matches exactly
# ---------------------------------------------------------------------------

def oklch(L, C, H):
    h = math.radians(H)
    a, b = C * math.cos(h), C * math.sin(h)
    l_ = L + 0.3963377774 * a + 0.2158037573 * b
    m_ = L - 0.1055613458 * a - 0.0638541728 * b
    s_ = L - 0.0894841775 * a - 1.2914855480 * b
    l, m, s = l_ ** 3, m_ ** 3, s_ ** 3
    R = 4.0767416621 * l - 3.3077115913 * m + 0.2309699292 * s
    G = -1.2684380046 * l + 2.6097574011 * m - 0.3413193965 * s
    B = -0.0041960863 * l - 0.7034186147 * m + 1.7076147010 * s

    def g(x):
        x = max(0.0, min(1.0, x))
        return 12.92 * x if x <= 0.0031308 else 1.055 * (x ** (1 / 2.4)) - 0.055

    return RGBColor(round(g(R) * 255), round(g(G) * 255), round(g(B) * 255))

PAPER      = oklch(0.974, 0.009, 84)
PAPER_2    = oklch(0.948, 0.013, 82)
PAPER_3    = oklch(0.922, 0.015, 80)
INK        = oklch(0.245, 0.012, 62)
INK_2      = oklch(0.420, 0.012, 62)
INK_3      = oklch(0.580, 0.010, 62)
LINE       = oklch(0.875, 0.013, 80)
LINE_2     = oklch(0.800, 0.013, 80)
HUMAN      = oklch(0.585, 0.088, 52)    # clay / ochre
MACHINE    = oklch(0.555, 0.072, 252)   # slate-blue
INK_PANEL  = oklch(0.262, 0.013, 64)
PAPER_INK  = oklch(0.950, 0.010, 84)    # paper-on-ink
BRAND_RED  = RGBColor(0xFF, 0x00, 0x15)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
# muted text on dark backgrounds
INK_MUT    = oklch(0.66, 0.010, 70)
SLATE_SOFT = oklch(0.66, 0.050, 252)

# Brand fonts are Spectral / IBM Plex (Google Fonts). When those aren't
# installed, the site CSS falls back to Georgia / system-ui / ui-monospace —
# which on macOS resolve to the fonts below. Using them keeps the deck on-brand
# and fully portable (no font install needed). Swap these three back to
# "Spectral" / "IBM Plex Sans" / "IBM Plex Mono" if those fonts are installed.
SERIF = "Georgia"
SANS  = "Helvetica Neue"
MONO  = "Menlo"

EMU_IN = 914400
SW = Inches(13.333)
SH = Inches(7.5)
MX = Inches(0.92)          # left/right margin
CW = SW - 2 * MX           # content width

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

# ---------------------------------------------------------------------------
# Primitive helpers
# ---------------------------------------------------------------------------

def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    # send to back (it's already first, but be explicit)
    return s

def _set_tracking(run, pts):
    run._r.get_or_add_rPr().set("spc", str(int(pts * 100)))

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def hline(s, x, y, w, color=LINE, weight=1.0):
    ln = s.shapes.add_connector(2, x, y, x + w, y)  # 2 = straight
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln

def _apply_run(run, text, font, size, color, bold, italic, tracking):
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    # make Office honor the latin font for ascii
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn("a:latin"))
    if latin is None:
        latin = rPr.makeelement(qn("a:latin"), {})
        rPr.append(latin)
    latin.set("typeface", font)
    if tracking:
        _set_tracking(run, tracking)

def text(s, x, y, w, h, body, font=SANS, size=18, color=INK, bold=False,
         italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line=1.12, tracking=None, space_after=6):
    """body: str (\\n = new paragraph) OR list of paragraphs, where each
    paragraph is a str or a list of (text, kwargs) run tuples."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)

    if isinstance(body, str):
        paras = [body.split("\n")] if False else body.split("\n")
        paras = [p for p in body.split("\n")]
    else:
        paras = body

    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        if isinstance(para, str):
            runs = [(para, {})]
        else:
            runs = para
        for rtext, opts in runs:
            run = p.add_run()
            _apply_run(
                run, rtext,
                opts.get("font", font),
                opts.get("size", size),
                opts.get("color", color),
                opts.get("bold", bold),
                opts.get("italic", italic),
                opts.get("tracking", tracking),
            )
    return tb

def kicker(s, txt, color=INK_3, x=MX, y=Inches(0.62)):
    return text(s, x, y, Inches(8), Inches(0.4), txt.upper(),
                font=MONO, size=11.5, color=color, tracking=1.7)

def logomark(s, x, y, size):
    circ = rect(s, x, y, size, size, fill=BRAND_RED, shape=MSO_SHAPE.OVAL)
    inset = int(size * 0.12)
    star = rect(s, x + inset, y + inset, size - 2 * inset, size - 2 * inset,
                fill=WHITE, shape=MSO_SHAPE.STAR_4_POINT)
    return circ, star

def footer(s, idx, accent=HUMAN, dark=False):
    base = INK_MUT if dark else INK_3
    y = Inches(7.06)
    hline(s, MX, y, CW, color=(LINE_2 if not dark else INK_PANEL), weight=0.75)
    text(s, MX, y + Inches(0.06), Inches(4), Inches(0.3),
         [[("ORIENT", {"color": (PAPER_INK if dark else INK), "tracking": 2.0}),
           ("  ·  meaning OS", {"color": base, "tracking": 1.2})]],
         font=MONO, size=9)
    text(s, SW - MX - Inches(4), y + Inches(0.06), Inches(4), Inches(0.3),
         [[(f"{idx:02d}", {"color": accent}), (f" / 18", {"color": base})]],
         font=MONO, size=9, align=PP_ALIGN.RIGHT, tracking=1.2)

def image_card(s, path, x, y, w, accent=HUMAN, max_h=None):
    """Place an image (auto aspect) inside a subtle bordered card."""
    pic = s.shapes.add_picture(path, x, y, width=w)
    if max_h and pic.height > max_h:
        # rescale by height
        ratio = pic.width / pic.height
        pic.height = max_h
        pic.width = int(max_h * ratio)
        pic.left = x
    pic.line.color.rgb = LINE_2
    pic.line.width = Pt(1)
    pic.shadow.inherit = False
    return pic

# ---------------------------------------------------------------------------
# Reusable composite blocks
# ---------------------------------------------------------------------------

def heading(s, txt, x=MX, y=Inches(1.42), w=None, size=40, color=INK,
            runs=None):
    w = w or Inches(11.0)
    if runs is not None:
        return text(s, x, y, w, Inches(2.4), [runs], font=SERIF, size=size,
                    color=color, line=1.05, space_after=0)
    return text(s, x, y, w, Inches(2.4), txt, font=SERIF, size=size,
                color=color, line=1.05, space_after=0)

def two_cards(s, top, items, y=Inches(3.4), accent=HUMAN, dark=False):
    """items: list of (label, label_color, body_lines:list[str])"""
    gap = Inches(0.45)
    cw = (CW - gap) / 2
    fill = PAPER_2 if not dark else INK_PANEL
    edge = LINE if not dark else None
    body_color = INK_2 if not dark else INK_MUT
    for i, (label, lcolor, lines) in enumerate(items):
        x = MX + i * (cw + gap)
        rect(s, x, y, cw, Inches(2.7), fill=fill, line=edge, line_w=1.0)
        pad = Inches(0.34)
        text(s, x + pad, y + Inches(0.3), cw - 2 * pad, Inches(0.4),
             label.upper(), font=MONO, size=11, color=lcolor, tracking=1.5)
        hline(s, x + pad, y + Inches(0.74), cw - 2 * pad,
              color=(LINE_2 if not dark else INK_3), weight=0.75)
        text(s, x + pad, y + Inches(0.95), cw - 2 * pad, Inches(1.6),
             "\n".join(lines), font=SANS, size=15, color=body_color,
             line=1.22, space_after=7)

# ===========================================================================
# SLIDES
# ===========================================================================

# --- 01 · Title -------------------------------------------------------------
s = slide(INK)
logomark(s, Inches(0.92), Inches(0.86), Inches(0.62))
text(s, Inches(1.66), Inches(0.86), Inches(4), Inches(0.62), "Orient",
     font=SERIF, size=30, color=PAPER_INK, anchor=MSO_ANCHOR.MIDDLE)
text(s, MX, Inches(2.7), Inches(11.3), Inches(0.5),
     "MEANING OPERATING SYSTEM", font=MONO, size=13, color=SLATE_SOFT,
     tracking=2.2)
heading(s, "", x=MX, y=Inches(3.18), size=52, color=PAPER_INK,
        runs=[("A meaning operating system\nfor humans ", {}),
              ("and", {"italic": True, "color": HUMAN}),
              (" machines.", {})])
text(s, MX, Inches(5.55), Inches(10.5), Inches(0.9),
     "Where knowledge becomes meaning — captured, resolved, framed, "
     "measured and remembered — so people and the agents beside them act "
     "from the same understanding.",
     font=SANS, size=16, color=INK_MUT, line=1.3)
text(s, MX, Inches(6.78), Inches(11), Inches(0.4),
     "Investor keynote  ·  Stockholm  ·  orientos.ai", font=MONO, size=10.5,
     color=INK_3, tracking=1.6)

# --- 02 · Hook --------------------------------------------------------------
s = slide(PAPER)
kicker(s, "The thesis")
heading(s, "", y=Inches(2.35), size=50,
        runs=[("Making the work was never the hard part.\n", {}),
              ("Trusting it", {"italic": True, "color": HUMAN}),
              (" is.", {})])
hline(s, MX, Inches(4.7), Inches(3.0), color=HUMAN, weight=2.0)
text(s, MX, Inches(4.95), Inches(10.6), Inches(1.0),
     "Anyone can generate a deck, a memo, an answer. Knowing it is true, "
     "clear and actually understood — that is the work that's left.",
     font=SERIF, size=21, color=INK_2, italic=True, line=1.3)
footer(s, 2)

# --- 03 · Problem -----------------------------------------------------------
s = slide(PAPER)
kicker(s, "The problem")
heading(s, "More communication than ever.\nLess shared understanding.",
        y=Inches(1.4), size=38)
text(s, MX, Inches(3.5), Inches(11.0), Inches(1.2),
     [[("Organisations produce more decks, briefs, docs, updates and "
        "AI-generated material than ever — ", {"color": INK_2}),
       ("and assume alignment follows. It rarely does.",
        {"color": INK, "bold": True})]],
     font=SANS, size=19, line=1.35)
text(s, MX, Inches(4.95), Inches(11.0), Inches(1.4),
     "As ideas move across teams, meetings, tools — and now agents — meaning "
     "shifts, intent fragments, and decisions drift. The output keeps "
     "accelerating. The shared understanding does not.",
     font=SANS, size=17, color=INK_2, line=1.4)
footer(s, 3)

# --- 04 · Why now -----------------------------------------------------------
s = slide(INK)
kicker(s, "Why now", color=SLATE_SOFT)
heading(s, "Two forces are colliding.", y=Inches(1.4), size=40,
        color=PAPER_INK)
two_cards(s, None, [
    ("Force 01 · infinite output",
     HUMAN,
     ["AI can now generate unlimited content — fast, fluent, and",
      "impossible to fully trust. Volume is no longer the constraint."]),
    ("Force 02 · agents that act",
     SLATE_SOFT,
     ["Agents are starting to act on company knowledge directly —",
      "and they guess what's true, current, and allowed to be said."]),
], y=Inches(2.7), dark=True)
text(s, MX, Inches(5.7), Inches(11.4), Inches(0.9),
     [[("The bottleneck moves from producing to trusting. The edge is no "
        "longer faster output — it's ", {"color": INK_MUT}),
       ("durable alignment.", {"color": PAPER_INK, "bold": True,
                               "font": SERIF, "italic": True})]],
     font=SANS, size=19, line=1.35)
footer(s, 4, accent=MACHINE, dark=True)

# --- 05 · The insight: two customers ---------------------------------------
s = slide(PAPER)
kicker(s, "The core insight")
heading(s, "Orient has two customers of meaning.", y=Inches(1.4), size=38)
two_cards(s, None, [
    ("For humans",
     HUMAN,
     ["A thinking space to understand, decide, explain and",
      "remember what matters.",
      "",
      "Humans use Orient to think."]),
    ("For machines",
     MACHINE,
     ["Structured meaning, evidence and constraints that",
      "agents can act on.",
      "",
      "Agents use Orient to understand."]),
], y=Inches(2.65))
text(s, MX, Inches(5.7), Inches(11.4), Inches(0.9),
     [[("Old roadmap: build tools so humans make better outputs.   ",
        {"color": INK_3}),
       ("New roadmap: build the meaning infrastructure humans and agents "
        "share.", {"color": INK, "bold": True})]],
     font=SANS, size=15.5, line=1.35)
footer(s, 5)

# --- 06 · Retrieval vs meaning ---------------------------------------------
s = slide(INK)
kicker(s, "The difference", color=SLATE_SOFT)
heading(s, "Retrieval hands an agent chunks.\nOrient hands it understanding.",
        y=Inches(1.4), size=36, color=PAPER_INK)
two_cards(s, None, [
    ("Raw retrieval returns",
     INK_MUT,
     ["→  The top-k most similar passages",
      "→  No sense of what is current or stale",
      "→  No record of what's contested",
      "→  No constraints on what may be said"]),
    ("Orient returns",
     SLATE_SOFT,
     ["→  The current answer and its grounding status",
      "→  Support, contradictions and weak spots",
      "→  What's resolved vs. still an open question",
      "→  Whether it's safe to act — or escalate"]),
], y=Inches(3.45), dark=True)
footer(s, 6, accent=MACHINE, dark=True)

# --- 07 · The human loop ----------------------------------------------------
s = slide(PAPER)
kicker(s, "For humans · the core loop", color=HUMAN)
heading(s, "From raw material to shared understanding.", y=Inches(1.4),
        size=34, w=Inches(11.4))
steps = [("01", "Capture", "Bring scattered material into one place."),
         ("02", "Resolve", "Turn a question into a defensible answer."),
         ("03", "Frame", "Shape understanding others can use."),
         ("04", "Measure", "Check it's clear and ready to move people."),
         ("05", "Remember", "Hold on to what mattered, and how it changed.")]
n = len(steps)
gap = Inches(0.28)
cw = (CW - gap * (n - 1)) / n
y0 = Inches(2.9)
for i, (num, title, desc) in enumerate(steps):
    x = MX + i * (cw + gap)
    rect(s, x, y0, cw, Inches(2.5), fill=PAPER_2, line=LINE)
    text(s, x + Inches(0.2), y0 + Inches(0.22), cw, Inches(0.4), num,
         font=MONO, size=12, color=HUMAN, tracking=1.0)
    text(s, x + Inches(0.2), y0 + Inches(0.66), cw - Inches(0.3), Inches(0.5),
         title, font=SERIF, size=19, color=INK)
    text(s, x + Inches(0.2), y0 + Inches(1.18), cw - Inches(0.34), Inches(1.2),
         desc, font=SANS, size=12.5, color=INK_2, line=1.25)
text(s, MX, y0 + Inches(2.7), Inches(11.4), Inches(0.4),
     "↻  Remember feeds back into Capture — understanding compounds.",
     font=MONO, size=11.5, color=INK_3, tracking=0.6)
footer(s, 7)

# --- 08 · Demo: Capture -----------------------------------------------------
def demo_slide(idx, num, title, blurb, shot, chips_label, chips,
               accent=HUMAN):
    s = slide(PAPER)
    kicker(s, f"For humans · {num} {title}", color=accent)
    text(s, MX, Inches(1.32), Inches(5.0), Inches(0.7), title,
         font=SERIF, size=34, color=INK)
    text(s, MX, Inches(2.25), Inches(4.7), Inches(2.3), blurb,
         font=SANS, size=14.5, color=INK_2, line=1.34)
    text(s, MX, Inches(4.75), Inches(4.7), Inches(0.35), chips_label.upper(),
         font=MONO, size=10, color=INK_3, tracking=1.4)
    # chip row(s)
    cx, cy = MX, Inches(5.12)
    maxx = MX + Inches(4.7)
    for c in chips:
        wapprox = Inches(0.32 + len(c) * 0.085)
        if cx + wapprox > maxx:
            cx = MX
            cy = cy + Inches(0.52)
        rect(s, cx, cy, wapprox, Inches(0.4), fill=PAPER_2, line=LINE)
        text(s, cx, cy, wapprox, Inches(0.4), c, font=MONO, size=10.5,
             color=INK_2, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx = cx + wapprox + Inches(0.12)
    image_card(s, shot, Inches(6.0), Inches(1.32), Inches(6.4), accent=accent,
               max_h=Inches(4.7))
    footer(s, idx, accent=accent)
    return s

demo_slide(
    8, "01 ·", "Capture",
    "Paste a link, upload a doc, drop a note or voice memo, pull from "
    "connected tools. Orient reads the material and extracts what matters — "
    "claims, evidence, contradictions, decisions and open questions. "
    "Capture isn't storage; it's the beginning of understanding.",
    os.path.join(SHOTS, "screenshot-2-web.png"),
    "Drop anything in",
    ["links", "documents", "notes", "voice", "web", "transcripts", "media"],
)

# --- 09 · Demo: Resolve -----------------------------------------------------
demo_slide(
    9, "02 ·", "Resolve",
    "Where a real question becomes a defensible answer. Orient works through "
    "the sources, drafts an answer, and shows what holds up — a live boundary "
    "between knowledge and uncertainty, not just a summary.",
    os.path.join(SHOTS, "screenshot-3-web.png"),
    "It separates",
    ["resolved", "open probes", "weak", "contested", "stale"],
)

# --- 10 · Demo: Frame -------------------------------------------------------
demo_slide(
    10, "03 ·", "Frame",
    "Once a question is resolved, Orient shapes it into something others can "
    "use — and carries the evidence with the message, so communication stays "
    "connected to meaning instead of becoming polished drift.",
    os.path.join(SHOTS, "screenshot-4-web.png"),
    "A resolved answer becomes",
    ["a deck", "a memo", "an article", "a briefing", "a narrative"],
)

# --- 11 · Measure + Impact --------------------------------------------------
s = slide(PAPER)
kicker(s, "For humans · evaluation & memory", color=HUMAN)
heading(s, "A Meaning Health Score — before it ships,\nand after it lands.",
        y=Inches(1.4), size=32, w=Inches(11.4))
two_cards(s, None, [
    ("04 · Measure",
     HUMAN,
     ["Orient reads the work like a sharp editor: clarity,",
      "evidence, coherence, audience fit. A score on whether",
      "the meaning is strong enough to move a decision —",
      "not a vanity grade."]),
    ("06 · Impact",
     HUMAN,
     ["After it ships, Orient follows it into the world:",
      "comprehension, reach, clarity, speed — and where",
      "the meaning frayed. Know whether the idea landed,",
      "not just whether it was sent."]),
], y=Inches(3.5))
footer(s, 11)

# --- 12 · The platform: agent surfaces -------------------------------------
s = slide(INK)
kicker(s, "For machines · the platform", color=SLATE_SOFT)
heading(s, "Five things every agent can ask of Orient.", y=Inches(1.32),
        size=32, color=PAPER_INK, w=Inches(7.2))
surfaces = [
    ("Context", "The right working memory — question, current answer, claims, weak spots."),
    ("Grounding", "Every claim tied to support, contradiction, provenance and freshness."),
    ("Policy", "Durable rules — tone, lexical contracts, taboo phrases, accepted risks."),
    ("Task Framing", "A request turned into outcome, audience, evidence standard, risk."),
    ("Evaluation", "Did output stay faithful to source, intent and the lexical contract?"),
]
y0 = Inches(2.55)
rh = Inches(0.78)
for i, (t, d) in enumerate(surfaces):
    y = y0 + i * (rh + Inches(0.06))
    text(s, MX, y, Inches(0.6), rh, f"0{i+1}", font=MONO, size=13,
         color=SLATE_SOFT, anchor=MSO_ANCHOR.MIDDLE)
    text(s, MX + Inches(0.7), y, Inches(2.4), rh, t, font=SERIF, size=18,
         color=PAPER_INK, anchor=MSO_ANCHOR.MIDDLE)
    text(s, MX + Inches(3.0), y, Inches(3.85), rh, d, font=SANS, size=12,
         color=INK_MUT, anchor=MSO_ANCHOR.MIDDLE, line=1.18)
# code card
cx = Inches(8.35)
rect(s, cx, Inches(2.5), Inches(4.05), Inches(3.6), fill=INK_PANEL, line=INK_3,
     line_w=0.75)
text(s, cx + Inches(0.3), Inches(2.74), Inches(3.5), Inches(0.3),
     "// what an agent gets back", font=MONO, size=10.5, color=INK_MUT)
code_lines = [
    [("context ", {"color": PAPER_INK}), ("= ", {"color": SLATE_SOFT}),
     ("assembleMeaningContext(…)", {"color": HUMAN})],
    [("", {})],
    [("current_answer   → ", {"color": INK_MUT}),
     ("\"enterprise-first\"", {"color": PAPER_INK})],
    [("grounding_status → ", {"color": INK_MUT}),
     ("\"resolved\"", {"color": PAPER_INK})],
    [("weak_spots       → ", {"color": INK_MUT}),
     ("[\"pricing ageing\"]", {"color": PAPER_INK})],
    [("must_not_claim   → ", {"color": INK_MUT}),
     ("[\"compliance\"]", {"color": PAPER_INK})],
    [("safe_to_act      → ", {"color": INK_MUT}),
     ("true", {"color": SLATE_SOFT})],
]
text(s, cx + Inches(0.3), Inches(3.2), Inches(3.6), Inches(2.8), code_lines,
     font=MONO, size=11.5, line=1.45, space_after=2)
footer(s, 12, accent=MACHINE, dark=True)

# --- 13 · Resolver ----------------------------------------------------------
s = slide(INK)
kicker(s, "For machines · the killer surface", color=SLATE_SOFT)
heading(s, "Resolver — company uncertainty as an agent work-queue.",
        y=Inches(1.32), size=30, color=PAPER_INK, w=Inches(11.4))
text(s, MX, Inches(2.25), Inches(11.4), Inches(0.7),
     "Not a knowledge base — a live state of inquiry. What's known, decided, "
     "still open, too weak to claim, and when to ask a human.",
     font=SANS, size=15, color=INK_MUT, line=1.3)
lanes = [
    ("Resolved", "Safe to use", SLATE_SOFT,
     ["Enterprise-first, regulated verticals", "Positioning: \u201cmeaning layer\u201d"]),
    ("Open probes", "The work queue", HUMAN,
     ["Mid-market willingness to pay?", "Demand from compliance teams?"]),
    ("Contested", "Escalate to a human", RGBColor(0xC8, 0x7A, 0x4A),
     ["Team comprehension analytics", "Data-residency commitments"]),
    ("Stale", "Refresh before reuse", INK_MUT,
     ["Q2 pricing assumptions", "Regulatory example (2023)"]),
]
gap = Inches(0.3)
lw = (CW - gap * 3) / 4
ly = Inches(3.25)
for i, (t, sub, col, items) in enumerate(lanes):
    x = MX + i * (lw + gap)
    rect(s, x, ly, lw, Inches(3.05), fill=INK_PANEL, line=INK_3, line_w=0.75)
    rect(s, x, ly, lw, Inches(0.09), fill=col)
    text(s, x + Inches(0.22), ly + Inches(0.28), lw, Inches(0.4), t,
         font=SERIF, size=16, color=PAPER_INK)
    text(s, x + Inches(0.22), ly + Inches(0.72), lw, Inches(0.3), sub.upper(),
         font=MONO, size=9, color=col, tracking=1.0)
    for j, it in enumerate(items):
        ity = ly + Inches(1.16) + j * Inches(0.72)
        rect(s, x + Inches(0.22), ity, lw - Inches(0.44), Inches(0.62),
             fill=INK, line=INK_3, line_w=0.5)
        text(s, x + Inches(0.34), ity, lw - Inches(0.66), Inches(0.62), it,
             font=SANS, size=10.5, color=PAPER_INK,
             anchor=MSO_ANCHOR.MIDDLE, line=1.1)
footer(s, 13, accent=MACHINE, dark=True)

# --- 14 · Architecture / moat ----------------------------------------------
s = slide(PAPER)
kicker(s, "The architecture & the moat")
heading(s, "One substrate. Five layers. Two kinds of mind.", y=Inches(1.4),
        size=34, w=Inches(11.4))
layers = [
    ("Human Surfaces", "Capture, Resolve, Frame, Measure, Remember", "HUMANS", HUMAN),
    ("Meaning Objects", "Claims, evidence, answers, frames, audiences", "SHARED", INK_2),
    ("Agent Context Layer", "Context, answer, briefing & memory packs", "MACHINES", MACHINE),
    ("Agent Runtime", "APIs, SDK, MCP, grounding & meaning validation", "MACHINES", MACHINE),
    ("Feedback Loop", "What was read, believed, understood — remembered", "COMPOUNDS", HUMAN),
]
ly = Inches(2.7)
lh = Inches(0.7)
for i, (name, desc, who, col) in enumerate(layers):
    y = ly + i * (lh + Inches(0.07))
    rect(s, MX, y, CW, lh, fill=PAPER_2, line=LINE)
    rect(s, MX, y, Inches(0.09), lh, fill=col)
    text(s, MX + Inches(0.35), y, Inches(3.6), lh, name, font=SERIF, size=16,
         color=INK, anchor=MSO_ANCHOR.MIDDLE)
    text(s, MX + Inches(4.1), y, Inches(5.6), lh, desc, font=SANS, size=12.5,
         color=INK_2, anchor=MSO_ANCHOR.MIDDLE)
    text(s, SW - MX - Inches(2.2), y, Inches(2.0), lh, who, font=MONO,
         size=10, color=col, tracking=1.2, align=PP_ALIGN.RIGHT,
         anchor=MSO_ANCHOR.MIDDLE)
footer(s, 14)

# --- 15 · Who it's for / GTM ------------------------------------------------
s = slide(PAPER)
kicker(s, "Who it's for")
heading(s, "Organisations operating in complexity.", y=Inches(1.4), size=34)
cases = ["Leadership & strategy", "Internal alignment", "Customer comms",
         "R&D & research", "Regulation & governance", "Cross-functional ops"]
gap = Inches(0.3)
cw = (CW - gap * 2) / 3
for i, c in enumerate(cases):
    col = i % 3
    row = i // 3
    x = MX + col * (cw + gap)
    y = Inches(2.7) + row * (Inches(0.95) + Inches(0.22))
    rect(s, x, y, cw, Inches(0.95), fill=PAPER_2, line=LINE)
    text(s, x + Inches(0.28), y, cw - Inches(0.5), Inches(0.95), c,
         font=SERIF, size=17, color=INK, anchor=MSO_ANCHOR.MIDDLE)
text(s, MX, Inches(5.85), Inches(11.4), Inches(0.8),
     [[("Dual GTM:  ", {"color": HUMAN, "bold": True, "font": MONO,
                        "size": 13, "tracking": 1.0}),
       ("humans adopt the thinking space; agents create platform lock-in on "
        "the same meaning layer.", {"color": INK_2, "size": 16})]],
     font=SANS, line=1.3)
footer(s, 15)

# --- 16 · What Orient becomes ----------------------------------------------
s = slide(INK)
kicker(s, "The category", color=SLATE_SOFT)
heading(s, "One thinking space, three scales of meaning.", y=Inches(1.4),
        size=34, color=PAPER_INK, w=Inches(11.4))
tiers = [
    ("For an individual", "A calm thinking environment",
     "Gather material, ask better questions, remember what matters."),
    ("For a team", "Shared intelligence",
     "A living map of what's known, believed, decided and still open."),
    ("For an organisation", "A meaning system",
     "Not a repo or dashboard — a system that turns knowledge into orientation."),
]
gap = Inches(0.4)
cw = (CW - gap * 2) / 3
for i, (scope, title, desc) in enumerate(tiers):
    x = MX + i * (cw + gap)
    rect(s, x, Inches(2.8), cw, Inches(2.7), fill=INK_PANEL, line=INK_3,
         line_w=0.75)
    pad = Inches(0.32)
    text(s, x + pad, Inches(3.08), cw - 2 * pad, Inches(0.3), scope.upper(),
         font=MONO, size=10, color=SLATE_SOFT, tracking=1.3)
    text(s, x + pad, Inches(3.5), cw - 2 * pad, Inches(0.8), title,
         font=SERIF, size=20, color=PAPER_INK, line=1.05)
    text(s, x + pad, Inches(4.45), cw - 2 * pad, Inches(1.0), desc,
         font=SANS, size=13, color=INK_MUT, line=1.3)
text(s, MX, Inches(6.05), Inches(11.4), Inches(0.5),
     "Not a file repository  ·  not a chat archive  ·  not a dashboard  ·  "
     "not a generic AI assistant.", font=MONO, size=11, color=INK_3,
     tracking=0.8)
footer(s, 16, accent=MACHINE, dark=True)

# --- 17 · Team + vision -----------------------------------------------------
s = slide(PAPER)
kicker(s, "Team & vision")
heading(s, "", y=Inches(1.7), size=40,
        runs=[("The next advantage won't come from\nproducing more. It will "
               "come from being ", {}),
              ("understood.", {"italic": True, "color": HUMAN})])
hline(s, MX, Inches(4.4), Inches(3.0), color=HUMAN, weight=2.0)
text(s, MX, Inches(4.7), Inches(11.2), Inches(1.6),
     "Built by a team spanning reasoning systems, language models, agentic "
     "architectures, optimisation and vision. Across every domain the pattern "
     "was identical: output accelerated, but shared understanding did not. "
     "Orient is our answer to that gap — infrastructure for alignment.",
     font=SANS, size=16, color=INK_2, line=1.4)
footer(s, 17)

# --- 18 · The ask -----------------------------------------------------------
s = slide(INK)
logomark(s, Inches(0.92), Inches(0.86), Inches(0.56))
text(s, Inches(1.62), Inches(0.86), Inches(4), Inches(0.56), "Orient",
     font=SERIF, size=26, color=PAPER_INK, anchor=MSO_ANCHOR.MIDDLE)
heading(s, "Work from the same understanding.", y=Inches(2.5), size=46,
        color=PAPER_INK, w=Inches(11.4))
text(s, MX, Inches(4.2), Inches(10.8), Inches(1.0),
     "Orient is shared memory and meaning infrastructure for the age of "
     "agents. We're raising to build the platform and onboard design "
     "partners in regulated, high-complexity organisations.",
     font=SANS, size=17, color=INK_MUT, line=1.4)
# ask chips
asks = ["The round  →  ____", "Use of funds  →  platform + GTM",
        "Looking for  →  partners in complexity"]
ax = MX
for a in asks:
    w = Inches(0.5 + len(a) * 0.092)
    rect(s, ax, Inches(5.7), w, Inches(0.5), fill=INK_PANEL, line=INK_3,
         line_w=0.75)
    text(s, ax, Inches(5.7), w, Inches(0.5), a, font=MONO, size=11,
         color=PAPER_INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    ax = ax + w + Inches(0.2)
text(s, MX, Inches(6.7), Inches(11), Inches(0.4),
     "mikael@orientos.ai  ·  Malmskillnadsgatan 32, Stockholm",
     font=MONO, size=11, color=SLATE_SOFT, tracking=1.0)

out = os.path.join(HERE, "Orient-Keynote.pptx")
prs.save(out)
print("saved:", out, "| slides:", len(prs.slides._sldIdLst))
